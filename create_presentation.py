"""
TnP Connect - PDF to PPTX Template Converter
PDF ka content template ke format mein dalta hai
"""

import copy
import io
import os
import sys

# Force UTF-8 output encoding for Windows terminal compatibility
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from lxml import etree

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


# ─────────────────────────────────────────────
# Helper: slide ko deep copy karna
# ─────────────────────────────────────────────
def duplicate_slide(prs, slide_index):
    """Template slide ko duplicate karta hai aur naya slide return karta hai."""
    template_slide = prs.slides[slide_index]
    
    # Slide layout lo
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Template slide ka XML copy karo (deep copy)
    template_xml = copy.deepcopy(template_slide._element)
    
    # Naye slide ka spTree (shapes tree) replace karo
    new_spTree = new_slide._element.spTree
    template_spTree = template_xml.spTree
    
    # Existing shapes clear karo
    for child in list(new_spTree):
        new_spTree.remove(child)
    
    # Template ke shapes copy karo
    for child in list(template_spTree):
        new_spTree.append(copy.deepcopy(child))
    
    return new_slide


def set_text_in_shape(shape, text, font_size=None, bold=None, color=None, align=None):
    """Shape ke text frame mein text set karta hai."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    
    # Sab paragraphs clear karo
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    
    # Pehla paragraph use karo
    p = tf.paragraphs[0]
    p.clear()
    
    if align:
        p.alignment = align
    
    run = p.add_run()
    run.text = text
    
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_image_to_slide(slide, image_path, left, top, width, height):
    """Slide mein image add karta hai."""
    pic = slide.shapes.add_picture(image_path, left, top, width, height)
    return pic


def add_textbox(slide, text, left, top, width, height, 
                font_size=12, bold=False, color=RGBColor(0x33, 0x33, 0x33),
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Slide mein naya text box add karta hai."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────
# Main Script
# ─────────────────────────────────────────────
print("="*60)
print("TnP Connect - PPTX Presentation Builder")
print("="*60)

# Paths
TEMPLATE_PATH = "BE Project Presentation Template.pptx"
PDF_PATH = "TnP_Connect Ppt (1).pdf"
OUTPUT_PATH = "TnP_Connect_Final_Presentation.pptx"
IMG_FOLDER = "pdf_images"

# PDF images already extracted hain - reuse karein
print(f"\n[1] PDF pages images ready: {len(os.listdir(IMG_FOLDER))} images")

# Template load karo
print(f"[2] Loading template: {TEMPLATE_PATH}")
prs = Presentation(TEMPLATE_PATH)

slide_width = prs.slide_width   # 10 inches
slide_height = prs.slide_height  # 7.5 inches

print(f"    Slide size: {slide_width/914400:.1f}\" x {slide_height/914400:.1f}\"")
print(f"    Template slides: {len(prs.slides)}")

# ─────────────────────────────────────────────
# SLIDE 1: Cover Slide - Template ka pehla slide customize karo
# ─────────────────────────────────────────────
print("\n[3] Building Slide 1 - Cover Slide...")

cover_slide = prs.slides[0]

# Template mein cover slide ke shapes ko update karo
for shape in cover_slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        
        # Title field
        if "*Title*" in text or "BE Project-I Presentation" in text:
            # Clear and set new content
            tf = shape.text_frame
            tf.clear()
            
            # Line 1: Course info
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = "BE Project-I Presentation"
            r1.font.size = Pt(13)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            # Line 2: Academic year
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = "A.Y. 2026-27 (Odd)"
            r2.font.size = Pt(11)
            r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            # Line 3: Title
            p3 = tf.add_paragraph()
            r3 = p3.add_run()
            r3.text = "TnP Connect"
            r3.font.size = Pt(22)
            r3.font.bold = True
            r3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            # Line 4: Subtitle
            p4 = tf.add_paragraph()
            r4 = p4.add_run()
            r4.text = "Unified Training & Placement Management Solution"
            r4.font.size = Pt(12)
            r4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            print("    ✓ Title updated")
        
        # Group members field
        elif "Group Member Details" in text or "Presented By" in text:
            if "Presented By" in text:
                tf = shape.text_frame
                tf.clear()
                
                p1 = tf.paragraphs[0]
                r1 = p1.add_run()
                r1.text = "Presented By:"
                r1.font.size = Pt(11)
                r1.font.bold = True
                r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                
                members = [
                    "1. Neev Jain (31)",
                    "2. Pratyush Pandey (34)",
                    "3. Ankush Sahu (45)"
                ]
                for member in members:
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = member
                    r.font.size = Pt(10)
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                
                print("    ✓ Group members updated")
        
        # Guide field
        elif "Name of Guide" in text:
            tf = shape.text_frame
            tf.clear()
            
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = "Name of Guide:"
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = "[Guide Name]"
            r2.font.size = Pt(10)
            r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            
            print("    ✓ Guide field updated")

print("    ✓ Cover slide complete!")

# ─────────────────────────────────────────────
# SLIDES 2-19: Content Slides
# Template ke slide 2 ka layout use karo
# PDF pages 2-19 ko content area mein embed karo
# ─────────────────────────────────────────────

# Template ka slide 2 banao - ye content template hai
# Hum isko use karenge har page ke liye

# PDF ke ek-ek page ke liye slide banao (Page 2 se 19)
pdf_pages_data = [
    # (page_num, title, subtitle)
    (2,  "Table of Contents", "TnP Connect - Unified Training & Placement Management Solution"),
    (3,  "Objectives", "TnP Connect - What we aim to achieve"),
    (4,  "Proposed System Architecture", "TnP Connect - System Design"),
    (5,  "Entity Relationship Model", "ER Diagram of TnP Connect"),
    (6,  "Modules of TnP Connect", "Student & Admin Features"),
    (7,  "Modules - Internships & Applications", "Broadcast Alert | CS&E (CS) | TCET"),
    (8,  "Modules - Sessions & Registration", "CS&E (CS) | TCET | Division A"),
    (9,  "Modules - Resources & Notes", "TnP Connect - Resources Module"),
    (10, "Modules - Contests", "CS&E (Cybersecurity) | TCET"),
    (11, "Modules - Notifications", "CS&E (Cybersecurity) | TCET"),
    (12, "TnP Connect UI - Student Dashboard", "UI of the System"),
    (13, "TnP Connect UI - Admin Dashboard", "UI of the System"),
    (14, "TnP Connect UI - AI Features", "UI of the System"),
    (15, "Core Implementation - JWT Authentication", "CS&E (Cybersecurity) | TCET"),
    (16, "Core Implementation - CRUD Operations", "CS&E (Cybersecurity) | TCET"),
    (17, "Core Implementation - Validators", "CS&E (Cybersecurity) | TCET"),
    (18, "Core Implementation - Responsive UI", "CS&E (Cybersecurity) | TCET"),
    (19, "Thank You", "TnP_Connect: Bridging the Gap Between Academia and Industry"),
]

print(f"\n[4] Building {len(pdf_pages_data)} content slides...")

# Template slide 2 ka layout
content_layout = prs.slides[1].slide_layout

for idx, (page_num, title, subtitle) in enumerate(pdf_pages_data):
    slide_num = idx + 2  # Slide numbers 2 to 19
    print(f"    Building Slide {slide_num}: {title[:40]}...")
    
    # Naya slide add karo content layout se
    new_slide = prs.slides.add_slide(content_layout)
    
    # ── TCET Logo (top right) ──
    logo_path = "template_images/slide2_img1.jpg"
    if os.path.exists(logo_path):
        add_image_to_slide(
            new_slide,
            logo_path,
            left=Inches(8.81),
            top=Inches(0.16),
            width=Inches(0.94),
            height=Inches(0.88)
        )
    
    # ── CS&E Logo (top left) ──
    cse_logo_path = "template_images/slide2_img2.png"
    if os.path.exists(cse_logo_path):
        add_image_to_slide(
            new_slide,
            cse_logo_path,
            left=Inches(0.30),
            top=Inches(0.28),
            width=Inches(0.94),
            height=Inches(0.75)
        )
    
    # ── Slide Title ──
    title_box = add_textbox(
        new_slide,
        title,
        left=Inches(1.40),
        top=Inches(0.18),
        width=Inches(7.20),
        height=Inches(0.65),
        font_size=16,
        bold=True,
        color=RGBColor(0x1A, 0x47, 0x8A),  # Dark blue - template color
        align=PP_ALIGN.LEFT
    )
    
    # ── Top separator line ──
    # (simple rectangle as line)
    from pptx.util import Pt as PtUtil
    line = new_slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(1.05),
        Inches(10), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1A, 0x47, 0x8A)
    line.line.fill.background()
    
    # ── PDF Page Image (main content area) ──
    img_path = os.path.join(IMG_FOLDER, f"page_{page_num:02d}.png")
    if os.path.exists(img_path):
        add_image_to_slide(
            new_slide,
            img_path,
            left=Inches(0.20),
            top=Inches(1.15),
            width=Inches(9.60),
            height=Inches(5.80)
        )
    
    # ── Bottom bar with slide number ──
    bottom_bar = new_slide.shapes.add_shape(
        1,
        Inches(0), Inches(7.30),
        Inches(10), Inches(0.20)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(0x1A, 0x47, 0x8A)
    bottom_bar.line.fill.background()
    
    # ── Slide number text ──
    slide_num_box = add_textbox(
        new_slide,
        str(slide_num),
        left=Inches(9.30),
        top=Inches(7.20),
        width=Inches(0.50),
        height=Inches(0.25),
        font_size=9,
        bold=False,
        color=RGBColor(0x66, 0x66, 0x66),
        align=PP_ALIGN.RIGHT
    )

print(f"\n[5] All {len(prs.slides)} slides built!")

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
print(f"\n[6] Saving presentation as: {OUTPUT_PATH}")
prs.save(OUTPUT_PATH)
print(f"\n{'='*60}")
print(f"✅ SUCCESS! Presentation saved: {OUTPUT_PATH}")
print(f"   Total slides: {len(prs.slides)}")
print(f"{'='*60}")
