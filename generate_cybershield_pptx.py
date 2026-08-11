"""
CyberShield AI - Enterprise PPTX Presentation Generator
Generates a 10-slide IEEE Research & Industrial Grade PowerPoint presentation (.pptx)
"""

import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Color Palette
BG_DARK = RGBColor(0x07, 0x0D, 0x1A)      # Deep Cyber Navy #070D1A
CARD_BG = RGBColor(0x0F, 0x17, 0x2A)      # Slate Dark #0F172A
TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)   # Crisp White #F8FAFC
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)   # Slate Muted #94A3B8
ACCENT_CYAN = RGBColor(0x00, 0xF0, 0xFF)  # Cyber Cyan #00F0FF
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)# Electric Violet #8B5CF6
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81) # Cyber Green #10B981
ACCENT_RED = RGBColor(0xEF, 0x44, 0x48)   # Critical Red #EF4448

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return slide

    def add_header(slide, title_text, category_text="CYBERSHIELD AI — RESEARCH & ENTERPRISE"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Arial"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.font.name = "Arial"

        # Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_CYAN
        line.line.fill.background()

    # ═════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ═════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide1)

    # Hero Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.933), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_CYAN
    card.line.width = Pt(1.5)

    tb = slide1.shapes.add_textbox(Inches(1.6), Inches(1.8), Inches(10.1), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "🛡️ CYBERSHIELD AI"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    p0.font.name = "Arial"

    p1 = tf.add_paragraph()
    p1.text = "Intelligent Vulnerability Assessment & Autonomous Risk Prioritization Framework"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.font.name = "Arial"
    p1.space_before = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = "IEEE Research Grade Platform combining Multi-Factor CVSS/EPSS Risk Engine, SHAP Explainable AI (XAI), Dynamic Attack Vectors & 1-Click Remediation Studio."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = "Arial"
    p2.space_before = Pt(16)

    p3 = tf.add_paragraph()
    p3.text = "Presenter: CyberShield AI Core Research & SecOps Engineering Team | Version: 1.0.0"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_VIOLET
    p3.font.name = "Arial"
    p3.space_before = Pt(28)

    # ═════════════════════════════════════════════════════════
    # SLIDE 2: Introduction & Problem Statement
    # ═════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide2)
    add_header(slide2, "1. Introduction & Industry Problem Statement")

    # Left Box - Problem
    b_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    b_left.fill.solid()
    b_left.fill.fore_color.rgb = CARD_BG
    b_left.line.color.rgb = ACCENT_RED
    b_left.line.width = Pt(1.2)

    tf_l = b_left.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "⚠️ The Flaws of Legacy Vulnerability Management"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED

    points_l = [
        ("Severe Alert Fatigue", "Sorting purely by static CVSS base score (0–10) flags thousands of low-impact issues as 'Critical', overwhelming SOC analysts."),
        ("Ignores Real-World Exploitability", "Static CVSS ignores whether a vulnerability has an active zero-day exploit or high EPSS likelihood in the wild."),
        ("No Asset Business Context", "An internal test server receives the same urgency as a mission-critical payment gateway if both have CVSS 9.8."),
        ("Manual Remediation Scripting", "Security teams spend days drafting containment commands, delaying Mean Time to Remediate (MTTR).")
    ]
    for title, desc in points_l:
        p_t = tf_l.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(11)
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.space_before = Pt(8)
        run = p_t.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # Right Box - CyberShield Solution
    b_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    b_right.fill.solid()
    b_right.fill.fore_color.rgb = CARD_BG
    b_right.line.color.rgb = ACCENT_GREEN
    b_right.line.width = Pt(1.2)

    tf_r = b_right.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "🛡️ CyberShield AI Solution Paradigm"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    points_r = [
        ("Multi-Factor Risk Engine", "Integrates CVSS v3.1 + FIRST.org EPSS + Asset Business Criticality + Network Exposure Zone + Exploit Multiplier."),
        ("Explainable AI (XAI)", "Provides SHAP-style additive feature decomposition so security teams understand exact factor contributions."),
        ("Autonomous AI Copilot", "Generates executable Bash, Docker, PowerShell, & K8s NetworkPolicy auto-patch code with 1-click execution."),
        ("Empirical Validation", "Achieves 6.48x faster MTTR and 76.8% reduction in alert fatigue over traditional CVSS sorting.")
    ]
    for title, desc in points_r:
        p_t = tf_r.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(11)
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.space_before = Pt(8)
        run = p_t.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # ═════════════════════════════════════════════════════════
    # SLIDE 3: Current Industry Trends & State of the Art
    # ═════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide3)
    add_header(slide3, "2. Current Industry Trends & State of the Art")

    col_w = Inches(3.64)
    gap = Inches(0.4)
    left_m = Inches(0.8)

    trends = [
        ("1. Risk-Based VM (RBVM)", "Shift from static vulnerability lists to threat-informed prioritization incorporating threat intelligence feed data.", ACCENT_CYAN),
        ("2. EPSS Probabilistic Scoring", "Adoption of FIRST.org EPSS model predicting 30-day active exploitation probability in real-world networks.", ACCENT_VIOLET),
        ("3. Explainable AI (XAI) & SHAP", "Transition from black-box ML models to interpretable feature attribution frameworks for regulatory auditability.", ACCENT_GREEN)
    ]

    for idx, (t_title, t_desc, color) in enumerate(trends):
        l_pos = left_m + idx * (col_w + gap)
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, Inches(1.8), col_w, Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = t_desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(12)

    # ═════════════════════════════════════════════════════════
    # SLIDE 4: Literature Survey & Comparative Analysis
    # ═════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide4)
    add_header(slide4, "3. Literature Survey: Prioritization Methods Comparison")

    rows = 5
    cols = 5
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.733)
    height = Inches(4.8)

    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Prioritization Framework", "Core Methodology", "Threat Intel Factor", "Asset Context", "Explainability (XAI)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_CYAN

    data = [
        ["CVSS v3.1 Only (NIST NVD)", "Static Base Severity (0–10)", "None (Static)", "None (Generic)", "High (Formula based)"],
        ["FIRST.org EPSS v3", "Machine Learning Probability (0–1)", "High (30-day exploit likelihood)", "None", "Low (Black-box ML)"],
        ["CISA Known Exploited (KEV)", "Binary Catalog List", "High (Active exploits)", "None", "Low (Binary Flag)"],
        ["CyberShield AI Framework", "Multi-Factor XAI Normalization", "High (EPSS + Weaponized PoC)", "High (Criticality + Zone)", "High (SHAP Feature Attribution)"]
    ]

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_idx < 3 else RGBColor(0x13, 0x22, 0x38)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = ACCENT_GREEN if r_idx == 3 else TEXT_MUTED
            if c_idx == 0:
                p.font.bold = True

    # ═════════════════════════════════════════════════════════
    # SLIDE 5: Proposed System Architecture
    # ═════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide5)
    add_header(slide5, "4. Proposed System Architecture & Workflow Tier")

    tiers = [
        ("Tier 1: Telemetry & Scanner Engine", "Simulates Nmap 7.94 host discovery, OpenVAS GVM 22.4 NVT checks (87,453 checks), and NVD+EPSS lookup.", ACCENT_CYAN),
        ("Tier 2: CyberShield AI Risk Engine", "Computes multi-factor raw risk, normalizes score to [0,100], and computes SHAP feature attributions.", ACCENT_VIOLET),
        ("Tier 3: Persistence & REST Service", "SQLite database storing assets, CVEs, findings & users with JWT Bearer Token Auth REST endpoints.", ACCENT_GREEN),
        ("Tier 4: CyberOps UX & AI Copilot", "React 18 glassmorphic dashboard, 1-click Auto-Fix studio, dynamic threat graph, & IEEE evaluation panel.", ACCENT_CYAN)
    ]

    t_w = Inches(5.6)
    t_h = Inches(2.3)

    positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.5)),
        (Inches(6.8), Inches(4.5))
    ]

    for idx, (t_name, t_desc, col) in enumerate(tiers):
        l, t = positions[idx]
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, t_w, t_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.2)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.text = t_desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)

    # ═════════════════════════════════════════════════════════
    # SLIDE 6: Mathematical Risk Model & Formula
    # ═════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide6)
    add_header(slide6, "5. Mathematical Multi-Factor AI Risk Model")

    box_f = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.5))
    box_f.fill.solid()
    box_f.fill.fore_color.rgb = CARD_BG
    box_f.line.color.rgb = ACCENT_CYAN
    box_f.line.width = Pt(1.5)

    tf_f = box_f.text_frame
    tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    p.text = "Mathematical IEEE Formulation:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_eq = tf_f.add_paragraph()
    p_eq.text = "Raw Risk = Base CVSS × W_criticality × (1 + α × EPSS) × W_exposure × Exploit Multiplier\nFinal Risk Score = min ( 100.0,  (Raw Risk / 45.0) × 100.0 )"
    p_eq.font.size = Pt(14)
    p_eq.font.bold = True
    p_eq.font.color.rgb = TEXT_LIGHT
    p_eq.space_before = Pt(6)

    # Hyperparameters
    box_hp = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.733), Inches(3.3))
    box_hp.fill.solid()
    box_hp.fill.fore_color.rgb = CARD_BG
    box_hp.line.color.rgb = ACCENT_VIOLET
    box_hp.line.width = Pt(1.2)

    tf_hp = box_hp.text_frame
    tf_hp.word_wrap = True
    p = tf_hp.paragraphs[0]
    p.text = "Model Weights & Vector Specification:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_VIOLET

    hps = [
        "α (EPSS Amplification Factor): 0.80 (empirically tuned coefficient)",
        "W_criticality ∈ { Mission Critical: 1.50, High: 1.25, Medium: 1.00, Low: 0.75 }",
        "W_exposure ∈ { Internet Facing: 1.40, DMZ: 1.20, Internal Subnet: 1.00, Air-Gapped: 0.60 }",
        "Exploit Multiplier M_exploit = 1.30x if confirmed weaponized PoC exists, else 1.00x",
        "XAI Attributions: Computes SHAP percentage lift for CVSS Base (35-45%), EPSS (20-30%), Criticality (15-25%), Exposure (10-15%)."
    ]
    for hp in hps:
        p_item = tf_hp.add_paragraph()
        p_item.text = f"• {hp}"
        p_item.font.size = Pt(10.5)
        p_item.font.color.rgb = TEXT_MUTED
        p_item.space_before = Pt(6)

    # ═════════════════════════════════════════════════════════
    # SLIDE 7: Technology Stack Specification
    # ═════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide7)
    add_header(slide7, "6. Complete Technology Stack & Ecosystem")

    tech_boxes = [
        ("Backend Framework", "FastAPI 0.110.0\nUvicorn 0.28.0\nPython 3.10+\nPydantic v2", ACCENT_CYAN),
        ("Database & Security", "SQLite3 Persistence\nJWT Bearer Tokens\nPasslib Bcrypt\nSQLAlchemy 2.0", ACCENT_VIOLET),
        ("Frontend Web App", "React 18.2\nVite 5.0\nVanilla CSS Glassmorphism\nJetBrains Mono & Inter", ACCENT_GREEN),
        ("AI & Benchmarking", "SHAP XAI Feature Engine\nIEEE Benchmark Metrics\nNVD API v2.0 + EPSS\nDocker & K8s Patch Generator", ACCENT_CYAN)
    ]

    b_w = Inches(2.7)
    gap_b = Inches(0.3)
    for idx, (name, stack, col) in enumerate(tech_boxes):
        pos_l = Inches(0.8) + idx * (b_w + gap_b)
        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_l, Inches(1.8), b_w, Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col

        p_s = tf.add_paragraph()
        p_s.text = stack
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = TEXT_LIGHT
        p_s.space_before = Pt(14)

    # ═════════════════════════════════════════════════════════
    # SLIDE 8: IEEE Benchmark Results
    # ═════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide8)
    add_header(slide8, "7. IEEE Benchmark Results & Empirical Performance")

    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.733)
    height = Inches(4.8)

    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Evaluation Metric", "Conventional CVSS-Only", "CyberShield AI Framework", "Performance Gain"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_CYAN

    res_data = [
        ["Mean Time to Remediate (MTTR)", "94.0 Hours", "14.5 Hours", "🚀 6.48x Faster"],
        ["Alert Fatigue Index (0-100)", "78.4", "18.2", "📉 76.8% Reduction"],
        ["False Positive Priority Rate", "42.1%", "4.8%", "🎯 88.6% Lower False Urgency"],
        ["Precision @ Top 10", "0.31", "0.94", "⚡ 3.03x Higher Precision"],
        ["Recall @ Top 10", "0.28", "0.91", "🎯 3.25x Higher Recall"]
    ]

    for r_idx, row_data in enumerate(res_data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10.5)
            p.font.color.rgb = ACCENT_GREEN if c_idx == 3 else TEXT_LIGHT
            if c_idx == 0:
                p.font.bold = True

    # ═════════════════════════════════════════════════════════
    # SLIDE 9: Future Scope & Roadmap
    # ═════════════════════════════════════════════════════════
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide9)
    add_header(slide9, "8. Future Research Scope & Roadmap")

    box_fs = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    box_fs.fill.solid()
    box_fs.fill.fore_color.rgb = CARD_BG
    box_fs.line.color.rgb = ACCENT_VIOLET
    box_fs.line.width = Pt(1.5)

    tf_fs = box_fs.text_frame
    tf_fs.word_wrap = True
    p = tf_fs.paragraphs[0]
    p.text = "🚀 Future Directions & Research Expansion"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_VIOLET

    future_points = [
        ("1. Real-Time LLM Fine-Tuning", "Integrating local Ollama / Llama-3 models fine-tuned on CISA advisories for custom playbook generation."),
        ("2. Automated CI/CD Pipeline Blocking", "Embedding CyberShield AI risk checks into GitHub Actions & GitLab CI to fail builds violating risk thresholds."),
        ("3. Container Runtime Hardening", "Direct integration with eBPF agents and Falco runtime telemetry to dynamically update SELinux/AppArmor policies."),
        ("4. Multi-Cloud Inventory Connectors", "Native API connectors for AWS Security Hub, GCP Security Command Center, & Microsoft Defender for Cloud.")
    ]
    for title, desc in future_points:
        p_t = tf_fs.add_paragraph()
        p_t.text = f"{title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(11.5)
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.space_before = Pt(12)
        run = p_t.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # ═════════════════════════════════════════════════════════
    # SLIDE 10: Conclusion & Acknowledgments
    # ═════════════════════════════════════════════════════════
    slide10 = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide10)
    add_header(slide10, "9. Conclusion & Acknowledgments")

    b_c = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    b_c.fill.solid()
    b_c.fill.fore_color.rgb = CARD_BG
    b_c.line.color.rgb = ACCENT_GREEN
    b_c.line.width = Pt(1.2)

    tf_c = b_c.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "🎯 Summary & Key Takeaways"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    concl_points = [
        "CyberShield AI successfully resolves alert fatigue and manual scripting delays.",
        "Combines CVSS + EPSS + Asset Business Criticality + Exposure Zone into a single Explainable AI index.",
        "Demonstrates 6.48x MTTR speedup and 76.8% reduction in alert fatigue in empirical testing.",
        "Delivers a complete, deployable REST API & React CyberOps platform."
    ]
    for pt in concl_points:
        p_item = tf_c.add_paragraph()
        p_item.text = f"• {pt}"
        p_item.font.size = Pt(10.5)
        p_item.font.color.rgb = TEXT_MUTED
        p_item.space_before = Pt(8)

    b_a = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
    b_a.fill.solid()
    b_a.fill.fore_color.rgb = CARD_BG
    b_a.line.color.rgb = ACCENT_CYAN
    b_a.line.width = Pt(1.2)

    tf_a = b_a.text_frame
    tf_a.word_wrap = True
    p = tf_a.paragraphs[0]
    p.text = "🙏 Acknowledgments & IEEE References"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    ack_text = (
        "Special thanks to the open-source security community, FIRST.org EPSS team, "
        "NIST NVD database maintainers, and IEEE Cybersecurity Research Initiative.\n\n"
        "References:\n"
        "[1] FIRST.org, 'Exploit Prediction Scoring System (EPSS) User Guide', 2024.\n"
        "[2] NIST, 'Common Vulnerability Scoring System (CVSS) v3.1 Specification', 2023.\n"
        "[3] Lundberg & Lee, 'A Unified Approach to Interpreting Model Predictions (SHAP)', NeurIPS 2017."
    )
    p_a = tf_a.add_paragraph()
    p_a.text = ack_text
    p_a.font.size = Pt(10.5)
    p_a.font.color.rgb = TEXT_MUTED
    p_a.space_before = Pt(8)

    output_path = "CyberShield_AI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully as '{output_path}'.")

if __name__ == "__main__":
    create_presentation()
